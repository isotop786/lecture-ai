from pptx import Presentation

def extract_text_from_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    text = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text.append(f"[Slide {slide_num}] {shape.text.strip()}")

    return "\n".join(text)
